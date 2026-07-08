"""
AFAC2026 PPT V3 重新生成（15 页，嵌入 7 张图）
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = r'd:\AFAC2026金融智能创新大赛\quantinsight-deploy'
ASSETS = os.path.join(ROOT, 'submission', '03_正式文档_WORD', '_assets')
OUTPUT = os.path.join(ROOT, 'QuantInsight_Pro_Pitch_Deck_V3.pptx')

# 16:9 = 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色
PRIMARY = RGBColor(0x1F, 0x77, 0xB4)
SUCCESS = RGBColor(0x2C, 0xA0, 0x2C)
WARNING = RGBColor(0xFF, 0x7F, 0x0E)
DANGER = RGBColor(0xD6, 0x27, 0x28)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
GRAY = RGBColor(0x88, 0x88, 0x88)


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei'
    return tb


def add_bullets(slide, left, top, width, height, items, font_size=14, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f'• {item}'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    elif width:
        return slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(path, left, top, height=height)
    return slide.shapes.add_picture(path, left, top)


def add_bg_rect(slide, color, left, top, width, height):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def page_title(slide, title, subtitle=None):
    """统一页面标题样式"""
    add_bg_rect(slide, PRIMARY, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    add_text(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
             title, font_size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.3),
                 subtitle, font_size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 底部 logo
    add_text(slide, Inches(11), Inches(7.1), Inches(2.2), Inches(0.3),
             'AFAC2026 · QuantInsight Pro', font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # ========================================
    # Slide 1: 封面
    # ========================================
    s = prs.slides.add_slide(blank)
    add_bg_rect(s, PRIMARY, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    add_text(s, Inches(0.5), Inches(2), Inches(12.3), Inches(0.6),
             'AFAC2026 金融智能创新大赛 · 初创组', font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.2),
             'QuantInsight Pro', font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
             'AI 驱动的另类数据量化投研平台', font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
             'SHAP 可解释 AI · 11.4 年回测 · 永字资管 POC', font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             '2026-07-08', font_size=14, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ========================================
    # Slide 2: 目录
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '目录 | Agenda', '8 大模块 · 5 分钟路演')
    add_bullets(s, Inches(1.5), Inches(1.5), Inches(10), Inches(5.5), [
        '1. 市场痛点：中小私募 0 AI 投研',
        '2. 解决方案：SHAP 可解释 AI + 11.4 年 POC',
        '3. 技术架构：17 因子 + 3 模态 + SHAP 归因',
        '4. 商业模式：4 客群 × 3 订阅 + 9 渠道',
        '5. 财务预测：5 年 30→620 客户，ARR 1.98亿 → 5.85亿',
        '6. 团队与顾问：4 创始 + 5 顾问 + 5 校 MOU',
        '7. 风险预案：5 大风险 + 3 套应对',
        '8. 愿景：让 1.5 万家中小私募 0 成本拥有 AI 投研',
    ], font_size=20)

    # ========================================
    # Slide 3: 痛点
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '1. 痛点 | 中小私募的 AI 投研困境')
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
             '中国资管规模 200 万亿，其中 1.5 万家中小私募（< 50 亿）面临：',
             font_size=18, color=DARK)
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), [
        '① 买不起：Wind + 优矿 + 团队 = 500 万/年，中小私募难以承担',
        '② 不会用：黑盒 AI 策略不可解释，监管与客户均不信任',
        '③ 没数据：舆情/资金流/政策等另类数据散落，无人整合',
        '④ 落不了：策略回测与实盘脱节，业绩归因混乱',
    ], font_size=20)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             '现状：90% 中小私募靠 Excel + 人工研判，年化跑赢基准者不足 20%',
             font_size=18, color=DANGER, bold=True)

    # ========================================
    # Slide 4: 解决方案（含商业模式 9 宫格）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '2. 解决方案 | SHAP 可解释 AI + 11.4 年 POC')
    img_path = os.path.join(ASSETS, '01_business_model_canvas.png')
    add_image(s, img_path, Inches(0.3), Inches(1.2), width=Inches(8.5))
    add_text(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.5),
             '3 大核心：', font_size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(9.0), Inches(2.2), Inches(4.0), Inches(5), [
        '透明：SHAP 归因',
        '可证：11.4 年回测',
        '可负担：2.4 万/年',
    ], font_size=18)

    # ========================================
    # Slide 5: 技术架构 - SHAP
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '3. 技术 | SHAP 可解释 AI（核心壁垒）')
    add_bullets(s, Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.5), [
        '17 因子库：价值/动量/质量/规模/波动/情绪',
        'SHAP 归因：每个预测的因子贡献度可解释',
        '监管友好：自动生成可审计报告',
        '3 模态融合：舆情 + 资金流 + 政策',
        'MIT 开源：回测引擎已开源 5 项',
        '21/21 单元测试 100% PASS',
    ], font_size=18)
    add_text(s, Inches(7.5), Inches(1.5), Inches(5.5), Inches(0.5),
             '5 大技术壁垒', font_size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(7.5), Inches(2.2), Inches(5.5), Inches(4.5), [
        '① SHAP（业内独家）',
        '② 11.4 年回测（最长）',
        '③ 17 因子（最全）',
        '④ 3 模态融合（最广）',
        '⑤ MIT 开源（最透）',
    ], font_size=18, color=SUCCESS)

    # ========================================
    # Slide 6: 11.4 年回测（T35 修正）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '3. 技术 | 11.4 年回测（T35 修正后）', 'HS300 8.56% / ZZ500 24.48% / CYB 11.55%')
    img_path = os.path.join(ASSETS, '04_backtest_curve.png')
    add_image(s, img_path, Inches(0.3), Inches(1.2), width=Inches(8.5))
    add_text(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.5),
             '3 指数 5 策略：', font_size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(9.0), Inches(2.2), Inches(4.0), Inches(5), [
        'HS300 多因子 8.56%',
        '超越基准 3.10pp',
        'ZZ500 多因子 24.48%',
        'CYB 多因子 11.55%',
        '修正前 19.22% → 修正后 8.56%',
        'T35 引擎已开源',
    ], font_size=16)

    # ========================================
    # Slide 7: 客户案例 - 永字资管
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '4. 客户案例 | 永字资管战略合作')
    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             '2026-05 永字资管（薛永再）签署战略合作，POC 试点 5 只产品',
             font_size=20, color=DARK)
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), [
        '✅ 5 只产品使用 QuantInsight Pro 多因子策略',
        '✅ 平均 11.4 年回测年化 8.56% (T35 修正)',
        '✅ 监管合规：SHAP 归因报告自动生成',
        '✅ 客户反馈：替代原 Wind+优矿，节省 70% 成本',
    ], font_size=20)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
             '永字资管法人薛永再：作为推荐单位，为本项目战略背书（非参赛队员）',
             font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 8: 商业模式 9 宫格（独立展示）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '4. 商业模式 | 9 宫格（Business Model Canvas）')
    img_path = os.path.join(ASSETS, '01_business_model_canvas.png')
    add_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

    # ========================================
    # Slide 9: 客户细分 × 订阅矩阵
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '4. 商业模式 | 4 客群 × 3 订阅 LTV 矩阵')
    img_path = os.path.join(ASSETS, '06_customer_subscription_matrix.png')
    add_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

    # ========================================
    # Slide 10: LTV/CAC 雷达
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '4. 商业模式 | LTV/CAC 雷达图（核心指标）')
    img_path = os.path.join(ASSETS, '02_ltv_cac_radar.png')
    add_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(8.5))
    add_text(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.5),
             '核心指标：', font_size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(9.0), Inches(2.2), Inches(4.0), Inches(5), [
        'LTV/CAC = 82.2',
        '（行业 3.0）',
        'NRR = 140%',
        '（行业 110%）',
        'Y3 毛利率 72%',
        '（行业 65%）',
        'CAC 回收 8 个月',
    ], font_size=18, color=SUCCESS)

    # ========================================
    # Slide 11: NRR 漏斗
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '4. 商业模式 | NRR 140% 漏斗')
    img_path = os.path.join(ASSETS, '03_nrr_funnel.png')
    add_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

    # ========================================
    # Slide 12: 客户增长曲线
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '5. 财务 | 5 年 30 → 620 客户路线图')
    img_path = os.path.join(ASSETS, '05_client_growth.png')
    add_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

    # ========================================
    # Slide 13: 团队架构
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '6. 团队 | 4 创始 + 5 顾问 + 5 校 MOU')
    img_path = os.path.join(ASSETS, '07_team_structure.png')
    add_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

    # ========================================
    # Slide 14: 风险预案
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '7. 风险预案 | 5 大风险 + 3 套应对')
    add_bullets(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), [
        '① 监管风险：AI 投顾持牌 → 已与永字合规团队合作，SHAP 自动审计报告',
        '② 数据合规：另类数据来源 → 已签 5 家授权，使用前 100% 脱敏',
        '③ 技术风险：开源回测引擎被 fork → MIT License，专利 + 软著 + 持续创新',
        '④ 市场风险：竞品（Wind AI） → 差异化 SHAP + 中小私募长尾',
        '⑤ 团队风险：4 人创业稳定性 → 期权池 35% + 顾问委员会 5 位 + 5 校 MOU',
    ], font_size=18)

    # ========================================
    # Slide 15: 愿景
    # ========================================
    s = prs.slides.add_slide(blank)
    add_bg_rect(s, PRIMARY, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    add_text(s, Inches(0.5), Inches(2), Inches(12.3), Inches(0.6),
             '愿景 | Vision', font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.2),
             '让 1.5 万家中小私募', font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.2),
             '0 成本拥有 AI 投研能力', font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
             'QuantInsight Pro 团队 · 冯亦根/王宇寒/官馨/梁理智',
             font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             '2026.07.08 · 等待您的回音',
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    sz = os.path.getsize(OUTPUT)
    print(f'✅ PPT V3 已生成: {os.path.basename(OUTPUT)}')
    print(f'   大小: {sz/1024/1024:.2f} MB')
    print(f'   页数: {len(prs.slides)}')


if __name__ == '__main__':
    main()
