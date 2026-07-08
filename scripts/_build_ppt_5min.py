"""
AFAC2026 5 分钟路演 PPT V1（11 页，嵌入 5-7 张图）
配套 5 分钟决赛路演视频
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = r'd:\AFAC2026金融智能创新大赛\quantinsight-deploy'
ASSETS = os.path.join(ROOT, 'submission', '03_正式文档_WORD', '_assets')
OUTPUT = os.path.join(ROOT, 'QuantInsight_Pro_Pitch_Deck_5min_V1.pptx')

# 16:9 = 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色（与 V3 一致）
PRIMARY = RGBColor(0x1F, 0x77, 0xB4)
SUCCESS = RGBColor(0x2C, 0xA0, 0x2C)
WARNING = RGBColor(0xFF, 0x7F, 0x0E)
DANGER = RGBColor(0xD6, 0x27, 0x28)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
GRAY = RGBColor(0x88, 0x88, 0x88)


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei'
    return tb


def add_bullets(slide, left, top, width, height, items, font_size=14, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f'• {item}'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    elif width:
        return slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(path, left, top, height=height)
    return slide.shapes.add_picture(path, left, top)


def add_bg_rect(slide, color, left, top, width, height):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def page_title(slide, title, subtitle=None):
    """统一页面标题样式（顶部蓝条 + 白字），底部统一页脚"""
    add_bg_rect(slide, PRIMARY, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    add_text(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
             title, font_size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.3),
                 subtitle, font_size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 底部页脚（统一）
    add_text(slide, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3),
             'AFAC2026 · 5 分钟路演  |  QuantInsight Pro',
             font_size=10, color=GRAY, align=PP_ALIGN.CENTER)


def safe_image(slide, path, left, top, width=None, height=None):
    """若文件存在则插入并返回 picture；否则返回 None（不抛错）"""
    if path and os.path.exists(path):
        try:
            return add_image(slide, path, left, top, width=width, height=height)
        except Exception as e:
            print(f'   ⚠️ 图片插入失败: {path} ({e})')
            return None
    print(f'   ⚠️ 图片缺失: {path}')
    return None


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]
    missing = []

    # ========================================
    # Slide 1: 封面（30s）
    # ========================================
    s = prs.slides.add_slide(blank)
    add_bg_rect(s, PRIMARY, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             'AFAC2026 金融智能创新大赛 · 决赛路演', font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.2),
             'QuantInsight Pro', font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.6),
             'AI 驱动的另类数据量化投研平台', font_size=26, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
             'SHAP 可解释 AI  ·  11.4 年回测  ·  永字资管 POC', font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
             '5 分钟路演', font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             '2026-07-08  ·  AFAC2026 · QuantInsight Pro',
             font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 2: Hook 30s —— 一个震撼数据
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, 'Hook | 90% 中小私募 跑不赢基准', '中小私募买不起 Wind · 90% 跑不赢基准')
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             '中国资管规模 200 万亿，其中 1.5 万家中小私募：',
             font_size=20, color=DARK)
    # 核心震撼数据
    add_bg_rect(s, DANGER, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.6))
    add_text(s, Inches(1.0), Inches(2.45), Inches(11.3), Inches(0.6),
             '中小私募 90% 跑不赢基准', font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(3.15), Inches(11.3), Inches(0.6),
             'Wind + 优矿 + 团队 = 500 万/年，中小私募根本买不起',
             font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             '结果：', font_size=22, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
        '① 90% 中小私募跑不赢基准 —— AI 投研是 0',
        '② 500 万/年的工具门槛 —— 99% 中小私募用 Excel + 人工研判',
        '③ 1.5 万家 × 50 万/年 —— 75 亿/年的 TAM，QuantInsight 切入',
    ], font_size=20, color=DARK)

    # ========================================
    # Slide 3: 痛点+方案 60s
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '痛点 & 方案 | 4 大痛点 + 3 大核心', '中小私募的 AI 投研困境 vs 我们的应对')
    add_text(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5),
             '4 大痛点：', font_size=22, bold=True, color=DANGER)
    add_bullets(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(5.0), [
        '① 买不起：Wind 500 万/年',
        '② 不会用：黑盒 AI 不可解释',
        '③ 没数据：另类数据散落',
        '④ 落不了：回测与实盘脱节',
    ], font_size=18, color=DARK)
    add_text(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.5),
             '3 大核心：', font_size=22, bold=True, color=SUCCESS)
    add_bullets(s, Inches(7.2), Inches(1.8), Inches(5.8), Inches(5.0), [
        '① 透明：SHAP 归因，监管友好',
        '② 可证：11.4 年 POC 回测',
        '③ 可负担：2.4 万/年（Wind 的 1/200）',
    ], font_size=18, color=DARK)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '现状 → 方案：QuantInsight Pro = SHAP 可解释 AI + 11.4 年回测 + 2.4 万/年',
             font_size=18, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 4: 技术 SHAP 60s（嵌入 02_ltv_cac_radar.png）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '技术 | SHAP 可解释 AI（核心壁垒）',
               '5 大壁垒 · 17 因子 · 3 模态融合 · MIT 开源')
    img_path = os.path.join(ASSETS, '02_ltv_cac_radar.png')
    pic = safe_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(7.5))
    if pic is None:
        missing.append(img_path)
    add_text(s, Inches(8.3), Inches(1.3), Inches(4.8), Inches(0.5),
             '5 大技术壁垒：', font_size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(8.3), Inches(1.9), Inches(4.8), Inches(5.0), [
        '① SHAP 归因（业内独家）',
        '② 11.4 年回测（业内最长）',
        '③ 17 因子库（业内最全）',
        '④ 3 模态融合（舆情/资金/政策）',
        '⑤ MIT 开源 + 软著',
    ], font_size=18, color=SUCCESS)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '21/21 单元测试 100% PASS  ·  客户信任：永字资管 POC 5 只产品',
             font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 5: 11.4 年回测 60s（嵌入 04_backtest_curve.png）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '回测 | 11.4 年 POC（T35 修正后）',
               'HS300 8.56% / ZZ500 24.48% / CYB 11.55%')
    img_path = os.path.join(ASSETS, '04_backtest_curve.png')
    pic = safe_image(s, img_path, Inches(0.3), Inches(1.2), width=Inches(8.5))
    if pic is None:
        missing.append(img_path)
    add_text(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.5),
             '3 指数 5 策略：', font_size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(9.0), Inches(2.1), Inches(4.0), Inches(5), [
        'HS300 多因子  8.56%',
        '超越基准  +3.10pp',
        'ZZ500 多因子 24.48%',
        'CYB  多因子  11.55%',
        '修正前 19.22% → 修正后 8.56%',
        'T35 引擎已 MIT 开源',
    ], font_size=16)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '最长时间窗口 + T35 修正 + MIT 开源 = 业内可信度第一',
             font_size=16, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 6: 客户案例 30s —— 永字资管 POC
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '客户案例 | 永字资管战略合作',
               '2026-05 签署  ·  POC 5 只产品  ·  节省 70% 成本')
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             '永字资管（薛永再 · 推荐单位 · 非参赛队员）',
             font_size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_bg_rect(s, SUCCESS, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.4))
    add_text(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.6),
             'POC 5 只产品全部使用 QuantInsight Pro 多因子策略',
             font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.6),
             '替代原 Wind+优矿  ·  平均年化 8.56%  ·  监管合规 SHAP',
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5), [
        '✅ 5 只产品：覆盖股票多头 / 量化对冲 / 主观选股',
        '✅ 节省 70% 工具成本：原 500 万/年 → 现 150 万/年',
        '✅ 监管友好：SHAP 归因报告自动生成，可审计',
        '✅ 法人薛永再战略背书（推荐单位）',
    ], font_size=18, color=DARK)

    # ========================================
    # Slide 7: 商业模式 60s（嵌入 01_business_model_canvas.png）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '商业模式 | 9 宫格（Business Model Canvas）',
               '4 客群 × 3 订阅 + 9 渠道 + 3 大护城河')
    img_path = os.path.join(ASSETS, '01_business_model_canvas.png')
    pic = safe_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))
    if pic is None:
        missing.append(img_path)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             'LTV/CAC = 82.2（行业 3.0）  ·  NRR = 140%  ·  Y3 毛利率 72%',
             font_size=16, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 8: 财务 60s（嵌入 05_client_growth.png）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '财务 | 5 年 30 → 620 客户',
               'ARR 1.98 亿 → 5.85 亿  ·  Y3 毛利率 72%')
    img_path = os.path.join(ASSETS, '05_client_growth.png')
    pic = safe_image(s, img_path, Inches(0.3), Inches(1.2), width=Inches(8.5))
    if pic is None:
        missing.append(img_path)
    add_text(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.5),
             '5 年财务：', font_size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(9.0), Inches(2.1), Inches(4.0), Inches(5), [
        'Y1: 30 客户  ·  ARR 1.98 亿',
        'Y3: 220 客户  ·  ARR 4.20 亿',
        'Y5: 620 客户  ·  ARR 5.85 亿',
        '毛利率:  62% → 72%',
        'CAC 回收: 8 个月',
        'LTV/CAC: 82.2',
    ], font_size=16)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             'TAM 75 亿 / SAM 22 亿 / SOM 5.6 亿（首期 5%）',
             font_size=16, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 9: 团队 30s（嵌入 07_team_structure.png）
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '团队 | 4 创始 + 5 顾问 + 5 校 MOU',
               '期权池 35%  ·  顾问委员会 5 位  ·  5 校人才 MOU')
    img_path = os.path.join(ASSETS, '07_team_structure.png')
    pic = safe_image(s, img_path, Inches(0.5), Inches(1.1), width=Inches(12.3))
    if pic is None:
        missing.append(img_path)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '冯亦根 / 王宇寒 / 官馨 / 梁理智  ·  复旦 / 清华 / 交大 / 上财 / 浙大',
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # ========================================
    # Slide 10: 风险 30s —— 5 大风险
    # ========================================
    s = prs.slides.add_slide(blank)
    page_title(s, '风险预案 | 5 大风险 + 3 套应对',
               '监管 / 数据 / 技术 / 市场 / 团队  ·  已全部闭环')
    add_bullets(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), [
        '① 监管风险：AI 投顾持牌 → 已与永字合规团队合作，SHAP 自动审计报告',
        '② 数据合规：另类数据来源 → 已签 5 家授权，使用前 100% 脱敏',
        '③ 技术风险：开源回测引擎被 fork → MIT License + 专利 + 软著 + 持续创新',
        '④ 市场风险：竞品（Wind AI） → 差异化 SHAP + 中小私募长尾 + 2.4 万/年低价',
        '⑤ 团队风险：4 人创业稳定性 → 期权池 35% + 顾问委员会 5 位 + 5 校 MOU',
    ], font_size=18)

    # ========================================
    # Slide 11: 愿景 30s
    # ========================================
    s = prs.slides.add_slide(blank)
    add_bg_rect(s, PRIMARY, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
             '愿景 | Vision', font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.2),
             '让 1.5 万家中小私募', font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.2),
             '0 成本拥有 AI 投研能力', font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
             'QuantInsight Pro 团队 · 冯亦根 / 王宇寒 / 官馨 / 梁理智',
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
             'SHAP 可解释  ·  11.4 年回测  ·  永字资管 POC',
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             'AFAC2026 · 5 分钟路演  ·  2026-07-08  ·  等待您的回音',
             font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    sz = os.path.getsize(OUTPUT)
    print(f'✅ PPT 5min V1 已生成: {os.path.basename(OUTPUT)}')
    print(f'   大小: {sz/1024/1024:.2f} MB')
    print(f'   页数: {len(prs.slides)}')
    if missing:
        print(f'   ⚠️ 缺失图片（已跳过，不影响生成）: {len(missing)} 张')
        for p in missing:
            print(f'      - {p}')


if __name__ == '__main__':
    main()
