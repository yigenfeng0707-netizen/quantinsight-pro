"""
AFAC2026 商业计划书 PPT（可编辑 pptx）
=====================================
按 BP 八章节逻辑 + 图文并茂，嵌入工作区已有图表 + 新渲染专业图。
输出：QuantInsight_Pro_商业计划书_BP_PPT_V1.pptx
"""
import os
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(r'd:\AFAC2026金融智能创新大赛\quantinsight-deploy')
ASSETS = ROOT / 'delivery' / '06_图表素材'
BP_CHARTS = ASSETS / 'bp_charts'
WORD_CHARTS = ROOT / 'submission' / '03_正式文档_WORD' / '_charts'
WORD_ASSETS = ROOT / 'submission' / '03_正式文档_WORD' / '_assets'
OUTPUT = ROOT / 'QuantInsight_Pro_商业计划书_BP_PPT_V1.pptx'
DELIVERY = ROOT / 'delivery' / '03_PPT' / 'QuantInsight_Pro_商业计划书_BP_PPT_V1.pptx'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRIMARY = RGBColor(0x1F, 0x77, 0xB4)
SUCCESS = RGBColor(0x2C, 0xA0, 0x2C)
WARNING = RGBColor(0xFF, 0x7F, 0x0E)
DANGER = RGBColor(0xD6, 0x27, 0x28)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
GRAY = RGBColor(0x88, 0x88, 0x88)
SOFT = RGBColor(0xF5, 0xF8, 0xFA)


def img(*candidates):
    for p in candidates:
        if p and Path(p).is_file():
            return str(p)
    return None


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei'
    return tb


def add_bullets(slide, left, top, width, height, items, font_size=15, color=DARK, spacing=5):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f'• {item}'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = 'Microsoft YaHei'
        p.space_after = Pt(spacing)
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    if not path:
        return None
    if width and height:
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(path, left, top, width=width)
    if height:
        return slide.shapes.add_picture(path, left, top, height=height)
    return slide.shapes.add_picture(path, left, top)


def add_bg_rect(slide, color, left, top, width, height):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_card(slide, left, top, width, height, fill=SOFT):
    return add_bg_rect(slide, fill, left, top, width, height)


def page_header(slide, title, subtitle=None):
    add_bg_rect(slide, PRIMARY, Inches(0), Inches(0), SLIDE_W, Inches(0.85))
    add_text(slide, Inches(0.45), Inches(0.12), Inches(12.2), Inches(0.55),
             title, font_size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.55), Inches(12.2), Inches(0.28),
                 subtitle, font_size=11, color=WHITE)
    add_text(slide, Inches(10.2), Inches(7.15), Inches(2.9), Inches(0.28),
             'AFAC2026 · 商业计划书 PPT', font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    cols = len(headers)
    nrows = len(rows) + 1
    table = slide.shapes.add_table(nrows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = 'Microsoft YaHei'
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = 'Microsoft YaHei'
                    run.font.color.rgb = DARK
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return table


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # resolve images
    canvas = img(ASSETS / '01_business_model_canvas.png', WORD_ASSETS / '01_business_model_canvas.png')
    ltv = img(ASSETS / '02_ltv_cac_radar.png', WORD_ASSETS / '02_ltv_cac_radar.png')
    nrr = img(ASSETS / '03_nrr_funnel.png', WORD_ASSETS / '03_nrr_funnel.png')
    backtest = img(ASSETS / '04_backtest_curve.png', WORD_ASSETS / '04_backtest_curve.png')
    growth = img(ASSETS / '05_client_growth.png', WORD_ASSETS / '05_client_growth.png')
    matrix = img(ASSETS / '06_customer_subscription_matrix.png', WORD_ASSETS / '06_customer_subscription_matrix.png')
    team = img(ASSETS / '07_team_structure.png', WORD_ASSETS / '07_team_structure.png')
    arch = img(WORD_CHARTS / 'chart_architecture.png')
    strategy = img(WORD_CHARTS / 'chart_strategy_comparison.png')
    tam = img(BP_CHARTS / '08_tam_sam_som.png')
    fin = img(BP_CHARTS / '09_revenue_profit.png')
    angel = img(BP_CHARTS / '10_angel_use_of_funds.png')
    comp = img(BP_CHARTS / '11_competition_radar.png')
    bt_table = img(BP_CHARTS / '12_backtest_table.png')
    road = img(BP_CHARTS / '13_milestone_roadmap.png')

    # ========== 1 封面 ==========
    s = prs.slides.add_slide(blank)
    add_bg_rect(s, PRIMARY, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.45),
             'AFAC2026 金融智能创新大赛 · 初创组', font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.0),
             'QuantInsight Pro', font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
             '商业计划书（PPT 版）', font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.45),
             'AI 驱动的另类数据量化投研平台 · SHAP 可解释性深度集成',
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.8),
             '项目编号 2026FINTECH-FINT-0093  |  Demo https://3blue1brownlab.cn\n'
             '推荐单位：杭州永字资产管理有限公司  |  编制：慧点资本量化研究部',
             font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.35),
             '2026 年 7 月  |  可编辑 PPTX', font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # ========== 2 目录 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '目录 | 商业计划书八章节', '图文并茂 · 逻辑闭环 · 可编辑提交')
    items = [
        ('01', '执行摘要', '一句话定位 + 关键里程碑'),
        ('02', '市场分析', 'TAM/SAM/SOM + 竞争格局'),
        ('03', '产品与技术', '8 模块 + SHAP + 11.4 年回测'),
        ('04', '商业模式', '9 宫格 + LTV/CAC + NRR'),
        ('05', '团队介绍', '4 创始 + 永字战略合作'),
        ('06', '财务预测', '5 年营收 + 融资路径'),
        ('07', '风险评估', '8 大风险 + 压力测试'),
        ('08', '落地案例', '永字资管 POC（≥1 案例）'),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i % 4
        row = i // 4
        left = Inches(0.5 + col * 3.15)
        top = Inches(1.3 + row * 2.7)
        add_card(s, left, top, Inches(3.0), Inches(2.3))
        add_text(s, left + Inches(0.15), top + Inches(0.25), Inches(2.7), Inches(0.5),
                 num, font_size=28, bold=True, color=PRIMARY)
        add_text(s, left + Inches(0.15), top + Inches(0.9), Inches(2.7), Inches(0.45),
                 title, font_size=18, bold=True, color=DARK)
        add_text(s, left + Inches(0.15), top + Inches(1.45), Inches(2.7), Inches(0.6),
                 desc, font_size=12, color=GRAY)

    # ========== 3 执行摘要 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '01 执行摘要', '让量化投资更智能，让另类数据更普惠，让 AI 决策可解释')
    add_card(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(1.35))
    add_text(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(1.1),
             'QuantInsight Pro 面向专业机构投资者的资管科技平台：另类数据 + 开源大模型微调/RAG 智能体 + '
             '开源回测引擎。核心差异化——业内首家将 SHAP 可解释性深度集成到 A 股智能选股，满足算法备案与客户沟通双重需求。',
             font_size=15, color=DARK)
    # KPI cards
    kpis = [
        ('30→620', '客户数（2026–2030）'),
        ('8.56%', 'HS300 多因子年化（T35）'),
        ('LTV/CAC 82.2', '单位经济模型'),
        ('永字 POC', '已签署战略合作'),
    ]
    for i, (v, lab) in enumerate(kpis):
        left = Inches(0.4 + i * 3.2)
        add_card(s, left, Inches(2.75), Inches(3.0), Inches(1.7), fill=WHITE)
        add_bg_rect(s, PRIMARY, left, Inches(2.75), Inches(0.12), Inches(1.7))
        add_text(s, left + Inches(0.25), Inches(2.95), Inches(2.6), Inches(0.7),
                 v, font_size=24, bold=True, color=PRIMARY)
        add_text(s, left + Inches(0.25), Inches(3.7), Inches(2.6), Inches(0.5),
                 lab, font_size=13, color=GRAY)
    add_bullets(s, Inches(0.5), Inches(4.7), Inches(12.2), Inches(2.2), [
        '目标客户：9 类机构（中小私募/理财子/券商资管等），长尾市场 1.5 万+ 家',
        '融资路径：5 轮 4.35 亿（天使→Pre-IPO），基准 IRR 28%，60 月启动 IPO',
        '产品形态：云原生 SaaS + 私有化 + 监管沙盒；Demo 已上线 https://3blue1brownlab.cn',
    ], font_size=15)

    # ========== 4 市场痛点 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '02 市场分析 · 痛点', '中国资管 200 万亿+，资管科技渗透率 < 5%')
    pains = [
        ('买不起', 'Wind+优矿+团队 ≈ 500万/年\n中小私募难以承担', DANGER),
        ('不会用', '黑盒 AI 不可解释\n监管与客户均不信任', WARNING),
        ('没数据', '舆情/资金流/政策散落\n无人整合另类数据', PRIMARY),
        ('落不了', '回测与实盘脱节\n业绩归因混乱', SUCCESS),
    ]
    for i, (t, d, c) in enumerate(pains):
        left = Inches(0.4 + i * 3.2)
        add_card(s, left, Inches(1.2), Inches(3.0), Inches(3.2))
        add_bg_rect(s, c, left, Inches(1.2), Inches(3.0), Inches(0.55))
        add_text(s, left + Inches(0.15), Inches(1.28), Inches(2.7), Inches(0.4),
                 t, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), Inches(2.0), Inches(2.6), Inches(2.0),
                 d, font_size=14, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
             '现状：90% 中小私募靠 Excel + 人工研判，年化跑赢基准者不足 20%',
             font_size=16, bold=True, color=DANGER)
    add_bullets(s, Inches(0.5), Inches(5.3), Inches(12.2), Inches(1.6), [
        '结构性蓝海：5000+ 中小私募，客单价 50–200 万/年',
        'AI 投研窗口：90% 机构无 AI 投研智能体，12–18 月技术窗口',
        '可解释刚需：证监会算法备案推动 SHAP 等可解释 AI 需求',
    ], font_size=14)

    # ========== 5 TAM ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '02 市场分析 · 市场规模', 'TAM / SAM / SOM')
    add_image(s, tam, Inches(0.3), Inches(1.1), width=Inches(8.2))
    add_text(s, Inches(8.7), Inches(1.3), Inches(4.2), Inches(0.4),
             '机会判断', font_size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(8.7), Inches(1.9), Inches(4.2), Inches(4.5), [
        'TAM：资管科技整体市场',
        'SAM：中小私募长尾可服务',
        'SOM：5 年可触达份额',
        '渗透率 < 5%，CAGR 30–40%',
        '竞品忽视长尾定价带',
        'SHAP 形成合规差异化',
    ], font_size=14)

    # ========== 6 竞争 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '02 市场分析 · 竞争格局', '相对 Wind / 同花顺 / 恒生的差异化')
    add_image(s, comp, Inches(0.2), Inches(1.0), width=Inches(7.5))
    add_table(s, Inches(7.6), Inches(1.3), Inches(5.3), Inches(4.8),
              ['对手', '优势', '我们的差异'],
              [
                  ['Wind', '品牌/数据全', 'AI+价格-50%'],
                  ['同花顺', '双覆盖', '中小私募+SHAP'],
                  ['恒生', '基础设施', '垂直场景灵活'],
                  ['QI Pro', 'AI+SHAP', '长尾蓝海'],
              ],
              col_widths=[Inches(1.3), Inches(1.8), Inches(2.2)])

    # ========== 7 产品方案 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '03 产品与技术 · 解决方案', '透明 · 可证 · 可负担')
    add_image(s, canvas, Inches(0.25), Inches(1.05), width=Inches(8.3))
    add_text(s, Inches(8.8), Inches(1.3), Inches(4.1), Inches(0.4),
             '三大核心', font_size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(8.8), Inches(1.9), Inches(4.1), Inches(4.8), [
        '透明：SHAP 因子归因',
        '可证：11.4 年真实回测',
        '可负担：2.4 万起/年',
        '8 大功能模块已上线',
        'Demo 生产环境可访问',
        'MIT 开源回测引擎',
    ], font_size=15)

    # ========== 8 八模块 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '03 产品与技术 · 八大功能模块', 'Demo 真实可运行')
    modules = [
        ('智能选股', '17 因子 Top10'),
        ('SHAP 解读', '业内独家归因'),
        ('AI 投研问答', 'RAG + 引用溯源'),
        ('量化回测', '5 策略 11.4 年'),
        ('智能盯盘', '7×24 异动预警'),
        ('模拟交易', 'A 股实时模拟'),
        ('自动报告', 'Word/PDF 导出'),
        ('实时看板', '北向资金热力'),
    ]
    for i, (t, d) in enumerate(modules):
        col, row = i % 4, i // 4
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.2 + row * 2.7)
        add_card(s, left, top, Inches(3.0), Inches(2.3))
        add_text(s, left + Inches(0.2), top + Inches(0.5), Inches(2.6), Inches(0.5),
                 f'{i+1}. {t}', font_size=18, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), top + Inches(1.2), Inches(2.6), Inches(0.6),
                 d, font_size=14, color=DARK, align=PP_ALIGN.CENTER)

    # ========== 9 技术架构 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '03 产品与技术 · 六层架构', '应用 → 智能体 → 服务 → 数据 → 算法 → 基础设施')
    if arch:
        add_image(s, arch, Inches(0.4), Inches(1.1), width=Inches(8.0))
    else:
        add_bullets(s, Inches(0.5), Inches(1.3), Inches(7.5), Inches(5), [
            '应用层：投研终端 / API / 可视化',
            '智能体层：投研助手 / 因子挖掘 / 自动报告',
            '服务层：AI / 策略 / 风控 / SHAP 解释器',
            '数据层：5 大数据源 + 17 因子库',
            '算法层：XGBoost + TreeExplainer + RAG',
            '基础设施：云原生 + 零信任安全',
        ], font_size=16)
    add_text(s, Inches(8.7), Inches(1.3), Inches(4.2), Inches(0.4),
             '技术壁垒', font_size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(8.7), Inches(1.9), Inches(4.2), Inches(4.8), [
        '① SHAP（业内独家）',
        '② 11.4 年回测（最长）',
        '③ 17 因子（最全）',
        '④ 3 模态融合（最广）',
        '⑤ MIT 开源（最透）',
        '21/21 单元测试 PASS',
    ], font_size=15, color=SUCCESS)

    # ========== 10 回测 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '03 产品与技术 · 11.4 年回测（T35 修正）',
                'HS300 8.56% / ZZ500 24.48% / CYB 11.55%')
    add_image(s, backtest, Inches(0.25), Inches(1.05), width=Inches(8.0))
    add_image(s, bt_table or strategy, Inches(8.3), Inches(1.2), width=Inches(4.7))
    add_text(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.4),
             '注：早期版本曾报告 19.22%，经 T35 开源引擎主动修正为 8.56%——不夸大也不掩饰。',
             font_size=12, color=GRAY)

    # ========== 11 商业模式 9宫格 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '04 商业模式 · Business Model Canvas', '9 宫格全景')
    add_image(s, canvas, Inches(0.4), Inches(1.0), width=Inches(12.5))

    # ========== 12 订阅矩阵 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '04 商业模式 · 客群 × 订阅', '4 客群 × 3 订阅 LTV 矩阵')
    add_image(s, matrix, Inches(0.4), Inches(1.0), width=Inches(12.5))

    # ========== 13 LTV/NRR ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '04 商业模式 · 单位经济', 'LTV/CAC = 82.2 · NRR = 140%')
    add_image(s, ltv, Inches(0.2), Inches(1.05), width=Inches(6.5))
    add_image(s, nrr, Inches(6.8), Inches(1.05), width=Inches(6.2))

    # ========== 14 收入结构 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '04 商业模式 · 收入与定价', 'SaaS 订阅为主（2030 占比 70%）')
    add_table(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(3.2),
              ['收入来源', '2030占比', '单价'],
              [
                  ['SaaS 订阅', '70%', '30–200万/年'],
                  ['定制开发', '15%', '100–500万/项目'],
                  ['数据 API', '7%', '按次计费'],
                  ['培训咨询', '4%', '5–20万/期'],
                  ['战略合作', '4%', '框架协议'],
              ])
    add_table(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(3.2),
              ['客户类型', '总价(万/年)', '典型模块'],
              [
                  ['中小私募', '50', '选股+SHAP+回测'],
                  ['银行理财子', '200', '全模块+风控'],
                  ['券商资管', '150', '投研+合规'],
                  ['高校研究所', '20', '学术版'],
                  ['战略合作', '框架', '试点+生态'],
              ])
    add_bullets(s, Inches(0.5), Inches(4.7), Inches(12.2), Inches(2.0), [
        '直销：2030 年销售 BD 80 人，覆盖银行/保险/公募',
        '渠道：券商/银行/律所联盟 + 清华/北大/上财学术渠道',
        '客户成功：目标续约率 92%，NRR 120%+',
    ], font_size=15)

    # ========== 15 团队 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '05 团队介绍', '4 创始 + 推荐单位战略背书')
    add_image(s, team, Inches(0.25), Inches(1.0), width=Inches(8.0))
    add_table(s, Inches(8.3), Inches(1.2), Inches(4.7), Inches(4.5),
              ['角色', '姓名'],
              [
                  ['CEO', '冯亦根'],
                  ['CTO', '王宇寒'],
                  ['产品/数据', '官馨'],
                  ['AI/量化', '梁理智'],
                  ['推荐单位', '薛永再(非队员)'],
              ],
              col_widths=[Inches(2.0), Inches(2.7)])
    add_text(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.4),
             '薛永再为杭州永字资管法定代表人 / 场外顾问，非参赛队员；战略合作已签署。',
             font_size=12, color=GRAY)

    # ========== 16 财务 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '06 财务预测 · 5 年营收', '基准场景：营收 0.20 → 7.95 亿')
    add_image(s, fin, Inches(0.25), Inches(1.0), width=Inches(8.2))
    add_image(s, growth, Inches(8.4), Inches(1.1), width=Inches(4.6))
    add_text(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.4),
             '2030 净利润 1.98 亿 · 净利率 25% · 5 年累计净利 2.94 亿',
             font_size=14, bold=True, color=PRIMARY)

    # ========== 17 融资 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '06 财务预测 · 融资与资金用途', '5 轮融资 4.35 亿 · 基准 IRR 28%')
    add_image(s, angel, Inches(0.3), Inches(1.1), width=Inches(6.0))
    add_table(s, Inches(6.5), Inches(1.2), Inches(6.4), Inches(4.5),
              ['轮次', '时点', '金额(万)', '投后估值(万)'],
              [
                  ['天使', '6月', '500', '5,000'],
                  ['Pre-A', '12月', '3,000', '30,000'],
                  ['A', '24月', '10,000', '80,000'],
                  ['B', '36月', '30,000', '200,000'],
                  ['Pre-IPO', '54月', '—', '500,000'],
              ])
    add_text(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.4),
             '退出：60 月启动 IPO（科创板/北交所）· IPO 估值 40–50 亿',
             font_size=14, color=DARK)

    # ========== 18 路线图 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '06 财务预测 · 客户拓展路线图', '种子 → 成长 → 扩张 → 领先 → 龙头')
    add_image(s, road or growth, Inches(0.5), Inches(1.2), width=Inches(12.3))

    # ========== 19 风险 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '07 风险评估', '8 大风险 + 压力测试')
    add_table(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(4.6),
              ['风险', '概率', '影响', '应对'],
              [
                  ['AI 技术迭代滞后', '20%', '-15%营收', 'R&D 15%营收'],
                  ['大客户流失', '15%', '-10%营收', '客户成功2x+长约'],
                  ['监管收紧', '30%', '-8%营收', '合规+监管沙盒'],
                  ['融资环境恶化', '25%', '-5pp IRR', '18月现金跑道'],
                  ['核心团队流失', '15%', '-5%营收', '期权池15%'],
                  ['竞品价格战', '40%', '-20%营收', 'SHAP差异化'],
                  ['数据源断供', '10%', '-25%营收', '多源备份'],
                  ['宏观衰退', '25%', '-15%营收', '9类客户分散'],
              ])
    add_text(s, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.8),
             '压力测试（2030）：保守 IRR 15% / 基准 28% / 乐观 38%  ·  缓释：技术护城河 + 客户分散 + 合规先行 + 人才绑定',
             font_size=14, color=DARK)

    # ========== 20 落地案例 ==========
    s = prs.slides.add_slide(blank)
    page_header(s, '08 落地案例 · 永字资管战略合作', '满足初创组「≥1 个落地客户/试点」要求')
    add_card(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(1.2))
    add_text(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.9),
             '合作方：杭州永字资产管理有限公司（中基协备案私募）  |  法定代表人：薛永再（推荐单位 / 场外顾问）\n'
             '状态：战略合作已签署  |  范围：AI 投研问答 + 量化回测 + SHAP 可解释选股',
             font_size=14, color=DARK)
    add_table(s, Inches(0.4), Inches(2.6), Inches(12.5), Inches(3.2),
              ['模块', '合作内容', '验证结果'],
              [
                  ['AI 投研问答', '财报/公告自然语言解读', '中文金融理解准确率 80%+'],
                  ['量化回测', 'HS300 多因子 11.4 年', '年化 8.56%，夏普 0.63'],
                  ['SHAP 解读', '17 因子决策归因', '满足算法备案可解释性'],
                  ['智能选股', '17 因子 Top10', '投研效率提升 65%（内部测算）'],
              ])
    add_text(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.6),
             '杭州私募证券基金管理人 471 家，86% 管理规模 < 10 亿 —— 精准服务长尾市场。',
             font_size=13, color=GRAY)

    # ========== 21 Ask / 愿景 ==========
    s = prs.slides.add_slide(blank)
    add_bg_rect(s, PRIMARY, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '愿景 · Ask', font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0),
             '让 1.5 万家中小私募\n0 成本拥有 AI 投研能力',
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.2),
             '寻求：专家初筛入围 · 投融资对接 · 金融机构 POC 试点\n'
             '团队：冯亦根 / 王宇寒 / 官馨 / 梁理智\n'
             'Demo：https://3blue1brownlab.cn  |  2026.07',
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    DELIVERY.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(OUTPUT, DELIVERY)
    sz = OUTPUT.stat().st_size
    print(f'OK {OUTPUT.name}')
    print(f'   size={sz/1024/1024:.2f} MB  slides={len(prs.slides)}')
    print(f'   copied -> {DELIVERY}')


if __name__ == '__main__':
    main()
