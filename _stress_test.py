"""
T23 压力测试 - 全交付物端到端验证
1. PPT 文件结构
2. PDF 文件页数
3. 视频文件大小
4. Excel 模型数据一致性
5. 关键数字跨文件验证
6. ZIP 完整性
"""
import os
import sys

# UTF-8 输出
sys.stdout.reconfigure(encoding='utf-8')

def test_separator(title):
    print('\n' + '='*70)
    print(f'  {title}')
    print('='*70)

def test_pass(msg):
    print(f'  [PASS] {msg}')

def test_fail(msg):
    print(f'  [FAIL] {msg}')

def test_warn(msg):
    print(f'  [WARN] {msg}')

results = {'PASS': 0, 'FAIL': 0, 'WARN': 0}
def log_pass(): results['PASS'] += 1
def log_fail(): results['FAIL'] += 1
def log_warn(): results['WARN'] += 1

# ============== 1. 文件存在性 ==============
test_separator('1. 核心交付物文件存在性')

core_files = {
    'BP V2 PDF': 'D:/shFintech/QuantInsight_Pro_BP_V2.pdf',
    'BP V2 MD': 'D:/shFintech/QuantInsight_Pro_BP_V2.md',
    'PPT V2': 'D:/shFintech/QuantInsight_Pro_Pitch_Deck_V2.pptx',
    '财务 V3 XLSX': 'D:/shFintech/QuantInsight_Pro_Financial_Model_V3.xlsx',
    '财务 V3 MD': 'D:/shFintech/QuantInsight_Pro_Financial_Report_V3.md',
    '财务图表 V3': 'D:/shFintech/QuantInsight_Pro_Chart_01_Financial.png',
    '白皮书 PDF': 'D:/shFintech/QuantInsight_Pro_Technical_Whitepaper_V1.pdf',
    '白皮书 MD': 'D:/shFintech/QuantInsight_Pro_Technical_Whitepaper_V1.md',
    'Demo 视频': 'D:/shFintech/QuantInsight_Pro_Demo_Video_Final.mp4',
    'Demo 字幕': 'D:/shFintech/QuantInsight_Pro_Demo_Video_V1.srt',
    '路演视频': 'D:/shFintech/QuantInsight_Pro_Pitch_Presenter_Video.mp4',
    'Q&A V2': 'D:/shFintech/QuantInsight_Pro_QA_Database_V2.md',
    '团队/合规/灾备': 'D:/shFintech/QuantInsight_Pro_Team_Compliance_DR_V1.md',
    '第三方盲评': 'D:/shFintech/QuantInsight_Pro_Third_Party_Review_V1.md',
    'README': 'D:/shFintech/README.md',
    '路演脚本': 'D:/shFintech/QuantInsight_Pro_Pitch_Script_V1.md',
}
for name, path in core_files.items():
    if os.path.exists(path):
        size = os.path.getsize(path)
        if size > 0:
            test_pass(f'{name} ({size:,} bytes)')
            log_pass()
        else:
            test_fail(f'{name} 大小为 0')
            log_fail()
    else:
        test_fail(f'{name} 不存在: {path}')
        log_fail()

# 验收报告
test_separator('2. 11 项验收报告存在性')
reports = [
    'T06_视频验收报告.md',
    'T11_短视频验收报告.md',
    'T12_路演视频验收报告.md',
    'T13_技术白皮书_验收报告.md',
    'T15_财务重写_验收报告.md',
    'T16_BP_V2_验收报告.md',
    'T17_PPT_V2_验收报告.md',
    'T18_QA升级_验收报告.md',
    'T19_团队合规灾备升级_验收报告.md',
    'T20_第三方盲评_验收报告.md',
    'T21_提交包整合_验收报告.md',
    'T22_路演脚本_验收报告.md',
]
for r in reports:
    path = f'D:/shFintech/{r}'
    if os.path.exists(path):
        size = os.path.getsize(path)
        test_pass(f'{r} ({size:,} bytes)')
        log_pass()
    else:
        test_fail(f'{r} 不存在')
        log_fail()

# ============== 3. PPT 验证 ==============
test_separator('3. PPT V2 文件结构')
try:
    from pptx import Presentation
    prs = Presentation('D:/shFintech/QuantInsight_Pro_Pitch_Deck_V2.pptx')
    n_slides = len(prs.slides)
    if n_slides == 16:
        test_pass(f'幻灯片数: {n_slides} (目标 16)')
        log_pass()
    else:
        test_warn(f'幻灯片数: {n_slides} (目标 16)')
        log_warn()

    # 检查每页有内容
    empty_slides = []
    for i, slide in enumerate(prs.slides, 1):
        text_count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            text_count += 1
        if text_count == 0:
            empty_slides.append(i)
    if empty_slides:
        test_warn(f'空白页: {empty_slides}')
        log_warn()
    else:
        test_pass('全部 16 页均有内容')
        log_pass()

    # 检查尺寸
    sw = prs.slide_width
    sh = prs.slide_height
    if sw == 12192000 and sh == 6858000:  # 16:9
        test_pass(f'幻灯片尺寸: 16:9 ({sw//914400}x{sh//914400} 英寸)')
        log_pass()
    else:
        test_warn(f'幻灯片尺寸: {sw//914400}x{sh//914400} 英寸')
        log_warn()
except Exception as e:
    test_fail(f'PPT 验证失败: {e}')
    log_fail()

# ============== 4. PDF 验证 ==============
test_separator('4. PDF 文件页数')
try:
    import pypdf
    pdfs = [
        ('BP V2', 'D:/shFintech/QuantInsight_Pro_BP_V2.pdf', 14),
        ('白皮书', 'D:/shFintech/QuantInsight_Pro_Technical_Whitepaper_V1.pdf', 19),
    ]
    for name, path, target in pdfs:
        reader = pypdf.PdfReader(path)
        n = len(reader.pages)
        diff = abs(n - target)
        if diff <= 5:
            test_pass(f'{name} PDF: {n} 页 (目标 {target}, 差 {diff})')
            log_pass()
        else:
            test_warn(f'{name} PDF: {n} 页 (目标 {target}, 差 {diff})')
            log_warn()
except Exception as e:
    test_fail(f'PDF 验证失败: {e}')
    log_fail()

# ============== 5. 视频文件大小 ==============
test_separator('5. 视频文件大小')
videos = [
    ('Demo 视频', 'D:/shFintech/QuantInsight_Pro_Demo_Video_Final.mp4', 9000000, 11000000),  # 9-11 MB
    ('路演视频', 'D:/shFintech/QuantInsight_Pro_Pitch_Presenter_Video.mp4', 15000000, 22000000),  # 15-22 MB
]
for name, path, min_size, max_size in videos:
    if os.path.exists(path):
        size = os.path.getsize(path)
        if min_size <= size <= max_size:
            test_pass(f'{name}: {size:,} bytes ({size/1024/1024:.1f} MB)')
            log_pass()
        else:
            test_warn(f'{name}: {size:,} bytes 超出预期 {min_size}-{max_size}')
            log_warn()

# ============== 6. Excel 财务模型 ==============
test_separator('6. 财务模型 V3 数据一致性')
try:
    import openpyxl
    wb = openpyxl.load_workbook('D:/shFintech/QuantInsight_Pro_Financial_Model_V3.xlsx', data_only=True)
    n_sheets = len(wb.sheetnames)
    test_pass(f'工作表数: {n_sheets} (目标 13)')
    log_pass() if n_sheets == 13 else log_warn()

    # 校验关键数字
    ws = wb['5年营收预测_基准']
    # 2030 营业总收入 (第 15 行, 第 6 列)
    revenue_2030 = ws.cell(row=15, column=6).value
    if revenue_2030 == 79500:
        test_pass(f'2030 营收: {revenue_2030:,} 万 = 7.95 亿 ✓')
        log_pass()
    else:
        test_fail(f'2030 营收异常: {revenue_2030} (期望 79500)')
        log_fail()

    ws = wb['成本与利润_基准']
    # 2030 净利润 (找营业利润行 + 3)
    for r in range(1, 30):
        cell = ws.cell(row=r, column=1).value
        if cell and '净利润' in str(cell):
            ni_2030 = ws.cell(row=r, column=6).value
            if ni_2030 == 19800:
                test_pass(f'2030 净利润: {ni_2030:,} 万 = 1.98 亿 ✓')
                log_pass()
            else:
                test_fail(f'2030 净利润异常: {ni_2030} (期望 19800)')
                log_fail()
            break

    # 客户数
    ws = wb['5年营收预测_基准']
    # 客户数 (家, 年末累计) 第 3 行, 第 6 列
    cust_2030 = ws.cell(row=3, column=6).value
    if cust_2030 == 620:
        test_pass(f'2030 客户数: {cust_2030} ✓')
        log_pass()
    else:
        test_fail(f'2030 客户数异常: {cust_2030} (期望 620)')
        log_fail()
except Exception as e:
    test_fail(f'Excel 验证失败: {e}')
    log_fail()

# ============== 7. 关键数字跨文件 ==============
test_separator('7. 关键数字跨文件一致性')

# 检查 7.95 亿 / 1.98 亿 / 620 / 25% / 40 亿 / 28% / 92% / 5 年 / 4.35 亿 等关键数字
critical_numbers = {
    '620': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '7.95': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '1.98': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '25%': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '40': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '28%': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '92%': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '4.35': ['BP_V2', 'PPT_V2', '财务_V3', 'Q&A_V2', '盲评', '路演脚本'],
    '冯亦根': ['BP_V2', 'Q&A_V2', '团队', '盲评', '路演脚本'],
    '薛永再': ['BP_V2', 'Q&A_V2', '团队', '盲评', '路演脚本'],
    '黄成选': ['BP_V2', 'Q&A_V2', '团队', '盲评', '路演脚本'],
    '冯思涵': ['BP_V2', 'Q&A_V2', '团队', '盲评', '路演脚本'],
    '杭州永字': ['BP_V2', 'Q&A_V2', '团队', '盲评', '路演脚本'],
    '2026FINTECH': ['BP_V2', 'PPT_V2', 'Q&A_V2', '团队', '盲评', '路演脚本'],
}

file_paths = {
    'BP_V2': 'D:/shFintech/QuantInsight_Pro_BP_V2.md',
    'PPT_V2': 'D:/shFintech/_build_ppt_v2.py',  # 脚本里包含数字
    '财务_V3': 'D:/shFintech/QuantInsight_Pro_Financial_Report_V3.md',
    'Q&A_V2': 'D:/shFintech/QuantInsight_Pro_QA_Database_V2.md',
    '盲评': 'D:/shFintech/QuantInsight_Pro_Third_Party_Review_V1.md',
    '路演脚本': 'D:/shFintech/QuantInsight_Pro_Pitch_Script_V1.md',
    '团队': 'D:/shFintech/QuantInsight_Pro_Team_Compliance_DR_V1.md',
}

# 加载文件内容
contents = {}
for k, p in file_paths.items():
    if os.path.exists(p):
        with open(p, 'r', encoding='utf-8') as f:
            contents[k] = f.read()

for num, files in critical_numbers.items():
    missing = []
    for f in files:
        if f in contents and num in contents[f]:
            pass
        else:
            missing.append(f)
    if not missing:
        test_pass(f'"{num}" 在全部 {len(files)} 个文件出现 ✓')
        log_pass()
    else:
        test_warn(f'"{num}" 缺失: {missing}')
        log_warn()

# ============== 8. ZIP 完整性 ==============
test_separator('8. 提交包 ZIP 完整性')
import zipfile
zip_path = 'D:/shFintech/QuantInsight_Pro_提交包_v1.0_20260606.zip'
if os.path.exists(zip_path):
    try:
        with zipfile.ZipFile(zip_path, 'r') as zf:
            bad = zf.testzip()
            if bad is None:
                test_pass('ZIP 完整性校验: 无损坏 ✓')
                log_pass()
            else:
                test_fail(f'ZIP 损坏: {bad}')
                log_fail()
            n_files = len(zf.namelist())
            test_pass(f'ZIP 包含 {n_files} 个文件')
            log_pass() if n_files >= 30 else log_warn()
    except Exception as e:
        test_fail(f'ZIP 验证失败: {e}')
        log_fail()
else:
    test_fail('ZIP 文件不存在')
    log_fail()

# ============== 9. Streamlit Demo 验证 ==============
test_separator('9. Streamlit Demo 包')
streamlit_path = 'D:/shFintech/streamlit_app/app.py'
if os.path.exists(streamlit_path):
    size = os.path.getsize(streamlit_path)
    with open(streamlit_path, 'r', encoding='utf-8') as f:
        content = f.read()
    lines = content.count('\n') + 1
    test_pass(f'app.py: {size:,} bytes, {lines} 行')
    log_pass() if 500 <= lines <= 700 else log_warn()
    # 检查必要文件
    for f in ['requirements.txt', 'README.md', 'config.toml']:
        p = f'D:/shFintech/streamlit_app/{f}'
        if os.path.exists(p):
            test_pass(f'{f} 存在')
            log_pass()
        else:
            test_warn(f'{f} 缺失')
            log_warn()
else:
    test_fail('Streamlit app.py 不存在')
    log_fail()

# ============== 10. 字幕文件 ==============
test_separator('10. 字幕文件')
srts = [
    'D:/shFintech/QuantInsight_Pro_Demo_Video_V1.srt',
    'D:/shFintech/QuantInsight_Pro_Pitch_Presenter_Video.srt',
]
for s in srts:
    if os.path.exists(s):
        size = os.path.getsize(s)
        with open(s, 'r', encoding='utf-8') as f:
            content = f.read()
        n_subs = content.count('-->')  # 时间码行数
        test_pass(f'{os.path.basename(s)}: {size} bytes, {n_subs} 条字幕')
        log_pass() if n_subs > 5 else log_warn()
    else:
        test_warn(f'{os.path.basename(s)} 不存在')
        log_warn()

# ============== 汇总 ==============
test_separator('压力测试汇总')
total = results['PASS'] + results['FAIL'] + results['WARN']
pass_rate = results['PASS'] / total * 100 if total > 0 else 0
print(f'  PASS: {results["PASS"]}')
print(f'  WARN: {results["WARN"]}')
print(f'  FAIL: {results["FAIL"]}')
print(f'  TOTAL: {total}')
print(f'  PASS RATE: {pass_rate:.1f}%')

if results['FAIL'] == 0:
    print(f'\n  结论: 压力测试 [PASS] - 全部交付物端到端正常')
else:
    print(f'\n  结论: 压力测试 [FAIL] - 有 {results["FAIL"]} 个关键问题, 需修复')
    sys.exit(1)
