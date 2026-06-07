"""
T17 PPT 重写 - 16 页资管科技路演版
基于 V3 财务 + V2 BP 9 章结构
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# 配色
PRIMARY = RGBColor(0x1F, 0x4E, 0x78)  # 深蓝
SECONDARY = RGBColor(0x2E, 0x86, 0xAB)  # 中蓝
ACCENT = RGBColor(0xA2, 0x3B, 0x72)  # 紫红
HIGHLIGHT = RGBColor(0x5A, 0x4F, 0xCF)  # 蓝紫
WARNING = RGBColor(0xF1, 0x8F, 0x01)  # 橙
LIGHT = RGBColor(0xD6, 0xE9, 0xF8)  # 浅蓝
GRAY = RGBColor(0x99, 0x99, 0x99)  # 灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)

# 创建 16:9 演示
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]

def set_run(run, text, size=18, bold=False, color=BLACK, font='Microsoft YaHei'):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 中文字体设置
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.find(qn('a:ea'))
    if eastAsia is None:
        eastAsia = etree.SubElement(rPr, qn('a:ea'))
    eastAsia.set('typeface', 'Microsoft YaHei')

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return tb

def add_title_bar(slide, title, subtitle=None):
    """添加标题栏"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run(run, title, size=24, bold=True, color=WHITE)

    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.4),
                 subtitle, size=13, color=GRAY, bold=False)

def add_footer(slide, page_num, total=16):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(8), Inches(0.3),
             f'QuantInsight Pro · 慧点资本 · 杭州永字资管 · 2026FINTECH-FINT-0093',
             size=8, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.4), Inches(0.3),
             f'{page_num} / {total}', size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def add_box(slide, x, y, w, h, fill=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    return shape

def add_table(slide, x, y, w, h, data, header_fill=PRIMARY, header_color=WHITE,
              body_size=11, header_size=12, first_col_bold=True, highlight_rows=None):
    """data: list of list, first row is header"""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table
    highlight_rows = highlight_rows or []
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val) if val is not None else ''
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(header_size if r == 0 else body_size)
                    run.font.bold = (r == 0) or (first_col_bold and c == 0)
                    if r == 0:
                        run.font.color.rgb = header_color
                    elif r in highlight_rows:
                        run.font.color.rgb = ACCENT
                    else:
                        run.font.color.rgb = BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table

# ============== Slide 1: 封面 ==============
slide = prs.slides.add_slide(BLANK)
# 背景渐变 (用大色块)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.line.fill.background()
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY

# 副标题条
sub = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), SLIDE_W, Inches(0.1))
sub.line.fill.background()
sub.fill.solid()
sub.fill.fore_color.rgb = ACCENT

# 主标题
add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
         'QuantInsight Pro', size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
         'AI 驱动的另类数据量化投研平台', size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(0.5),
         '— 资管科技 (AssetTech) 子公司视角 —', size=18, color=WHITE, align=PP_ALIGN.CENTER)

# 信息
add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.4),
         '慧点资本 (InsightQuant) 量化研究部', size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.4), Inches(11.3), Inches(0.4),
         '推荐单位: 杭州永字资产管理有限公司', size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4),
         '项目编号: 2026FINTECH-FINT-0093', size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
         'FinTech@外滩金融科技大赛 · 资管科技赛道', size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
         '2026 年 6 月', size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============== Slide 2: 项目概述 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '一、项目概述', '让量化投资更智能, 让另类数据更普惠, 让资管科技更本土化')

# 左侧: 项目定义
add_text(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.4),
         '项目定义', size=18, bold=True, color=PRIMARY)
add_text(slide, Inches(0.4), Inches(1.85), Inches(6.2), Inches(2.5),
         'QuantInsight Pro 是慧点资本联合杭州永字资产管理有限公司打造的新一代资管科技 (AssetTech) 平台, 通过整合另类数据、自研 AI 投研智能体和量化策略引擎, 为 9 类机构投资者提供数据+模型+策略+合规的一站式智能投研解决方案。',
         size=13, color=BLACK)

add_text(slide, Inches(0.4), Inches(4.0), Inches(6.2), Inches(0.4),
         '项目主体', size=18, bold=True, color=PRIMARY)
add_text(slide, Inches(0.4), Inches(4.4), Inches(6.2), Inches(2.5),
         '杭州永字资产管理有限公司资管科技子公司 (拟设立, 慧点资本 100% 控股)\n推荐单位: 杭州永字资产管理有限公司\n法定代表人: 薛永再 (创始团队 COO/合规)\n监管: 中基协私募基金管理人',
         size=13, color=BLACK)

# 右侧: 4 大核心价值
add_text(slide, Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.4),
         '4 大核心价值', size=18, bold=True, color=PRIMARY)

values = [
    ('技术壁垒', '自研 AI 投研智能体 (业内 90% 不具备)'),
    ('数据优势', '5 大数据源 (财报/公告/新闻/产业链/舆情)'),
    ('产品形态', '云原生 SaaS + 私有化 + 监管沙盒'),
    ('本土化', '学术合作 (清华/北大/上财) + A 股合规'),
]
y0 = 1.9
for i, (k, v) in enumerate(values):
    add_box(slide, Inches(7.0), Inches(y0 + i*1.0), Inches(6.0), Inches(0.85), fill=LIGHT, line_color=PRIMARY)
    add_text(slide, Inches(7.2), Inches(y0 + i*1.0 + 0.05), Inches(1.5), Inches(0.4),
             f'#{i+1} {k}', size=14, bold=True, color=PRIMARY)
    add_text(slide, Inches(8.7), Inches(y0 + i*1.0 + 0.15), Inches(4.0), Inches(0.6),
             v, size=12, color=BLACK)

add_footer(slide, 2)

# ============== Slide 3: 市场机会 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '二、市场机会', '中国资管 200 万亿+, 资管科技渗透率 < 5%, 5 年 CAGR 30-40%')

# 三大数字
nums = [
    ('200 万亿+', '中国资管行业 2025'),
    ('< 5%', '资管科技渗透率'),
    ('1.5 万+', '9 类机构总数'),
]
for i, (n, label) in enumerate(nums):
    add_box(slide, Inches(0.4 + i*4.4), Inches(1.3), Inches(4.0), Inches(1.5), fill=PRIMARY)
    add_text(slide, Inches(0.4 + i*4.4), Inches(1.4), Inches(4.0), Inches(0.7),
             n, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4 + i*4.4), Inches(2.1), Inches(4.0), Inches(0.6),
             label, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 9 类客户分群
add_text(slide, Inches(0.4), Inches(3.0), Inches(12.5), Inches(0.4),
         '9 类客户分群 × 1.5 万+ 机构', size=18, bold=True, color=PRIMARY)

data = [
    ['客户类型', '机构总数', '目标客户', '单价 (万)', '5 年累计', '策略'],
    ['中小私募', '10,000+', '500', '50', '500', '直销+推荐'],
    ['银行理财子', '30+', '20', '200', '20', '高层 BD+试点'],
    ['券商资管', '60+', '40', '150', '40', '战略联盟+共建'],
    ['信托公司', '50+', '30', '100', '30', '行业会议+案例'],
    ['保险资管', '30+', '20', '300', '20', '监管沙盒+试点'],
    ['公募基金', '150+', '15', '500', '15', '招投标+高层'],
    ['高校研究所', '200+', '50', '20', '50', '学术合作'],
    ['上市公司 IR', '5,000+', '100', '30', '100', '路演+案例'],
    ['战略合作', '20+', '20', '100', '20', '内部协同'],
    ['合计 / 加权', '15,540+', '795', '79', '795', '差异化蓝海'],
]
add_table(slide, Inches(0.4), Inches(3.5), Inches(12.5), Inches(3.4), data,
          body_size=11, header_size=12, highlight_rows=[10])

add_footer(slide, 3)

# ============== Slide 4: 产品方案 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '三、产品方案', 'AI 投研智能体 + 5 大数据源 + 9 类客户全覆盖')

# 4 大模块
modules = [
    ('另类数据中心', '5 大数据源\n• 传统金融数据\n• 新闻舆情\n• 产业链数据\n• 消费行为\n• 卫星/物流', SECONDARY),
    ('AI 投研引擎', '核心壁垒\n• 金融大模型 (DeepSeek/Qwen 微调)\n• 知识图谱 (1000 万+ 实体)\n• 因子挖掘 (自动 Alpha)\n• 智能问答', ACCENT),
    ('量化策略平台', '5.4 年回测已验证\n• 策略回测 (多周期)\n• 策略生成 (AI 辅助)\n• 风险模型 (VaR/CVaR)\n• 绩效归因', HIGHLIGHT),
    ('智能风控系统', '合规先行\n• 实时预警 (毫秒级)\n• 合规检查 (反洗钱)\n• 仓位管理 (动态)\n• 监管报送 (一键)', WARNING),
]
y0 = 1.4
for i, (name, desc, color) in enumerate(modules):
    x = 0.4 + (i % 4) * 3.2
    add_box(slide, Inches(x), Inches(y0), Inches(3.0), Inches(0.6), fill=color)
    add_text(slide, Inches(x), Inches(y0 + 0.1), Inches(3.0), Inches(0.4),
             name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(x), Inches(y0 + 0.6), Inches(3.0), Inches(2.4),
            fill=WHITE, line_color=color)
    add_text(slide, Inches(x + 0.1), Inches(y0 + 0.7), Inches(2.8), Inches(2.2),
             desc, size=11, color=BLACK)

# 底部: 5 大创新
add_text(slide, Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.4),
         '5 大技术亮点 (vs 业内基准)', size=16, bold=True, color=PRIMARY)
innov = [
    'AI 大模型: 业内 90% 机构不具备 → 12-18 月领先',
    '另类数据: 5 类融合 (Wind 仅 1-2 类) → 5x 数据广度',
    '知识图谱: 动态更新 (千万实体) → 实时性',
    '强化学习: 自适应策略 → 普惠化',
    '合规引擎: 监管沙盒 → 本土合规',
]
for i, item in enumerate(innov):
    x = 0.4 + (i % 5) * 2.5
    add_box(slide, Inches(x), Inches(4.7), Inches(2.4), Inches(1.5), fill=LIGHT)
    add_text(slide, Inches(x + 0.05), Inches(4.8), Inches(2.3), Inches(1.3),
             item, size=10, color=BLACK)

add_footer(slide, 4)

# ============== Slide 5: 技术架构 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '四、技术架构', '云原生 + AI 智能体编排 + 零信任安全')

# 4 层架构
layers = [
    ('应用层', '投研终端 | 移动 App | API 接口 | 数据可视化', PRIMARY),
    ('智能体层', '投研助手 | 因子挖掘 | 智能问答 | 自动报告', ACCENT),
    ('服务层', 'AI 引擎 | 策略引擎 | 风控引擎 | 数据服务', SECONDARY),
    ('数据层', '5 大数据源 + 知识图谱 + 自研因子库', HIGHLIGHT),
]
y0 = 1.5
for i, (name, desc, color) in enumerate(layers):
    h = 0.9
    y = y0 + i * (h + 0.2)
    add_box(slide, Inches(0.4), Inches(y), Inches(1.5), Inches(h), fill=color)
    add_text(slide, Inches(0.4), Inches(y + 0.25), Inches(1.5), Inches(0.4),
             name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(1.9), Inches(y), Inches(11.0), Inches(h), fill=WHITE, line_color=color)
    add_text(slide, Inches(2.0), Inches(y + 0.25), Inches(10.8), Inches(0.5),
             desc, size=14, color=BLACK)

# 基础设施
add_box(slide, Inches(0.4), Inches(5.7), Inches(1.5), Inches(0.9), fill=GRAY)
add_text(slide, Inches(0.4), Inches(5.95), Inches(1.5), Inches(0.4),
         '基础设施', size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_box(slide, Inches(1.9), Inches(5.7), Inches(11.0), Inches(0.9), fill=WHITE, line_color=GRAY)
add_text(slide, Inches(2.0), Inches(5.95), Inches(10.8), Inches(0.5),
         '云原生 | 弹性算力 | 实时数据流 | 零信任安全 | 99.9% SLA', size=14, color=BLACK)

# 右侧: 技术指标
add_text(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.3),
         '核心技术指标: 数据 10TB/日 | 响应 < 100ms | 可用性 99.9% | 回测 30s/5.4y | 并发 1000+ | API 100 万次/日',
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_footer(slide, 5)

# ============== Slide 6: 技术亮点 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '五、技术亮点', 'vs Wind/同花顺/恒生/通联 4 大业内基准')

data = [
    ['指标', 'QuantInsight Pro', 'Wind (万得)', '同花顺 iFinD', '恒生电子', '通联数据', '优势'],
    ['AI 投研智能体', '✓', '✗', '✗', '✗', '✗', '12-18 月领先'],
    ['5 类另类数据', '✓ (5 类)', '1-2 类', '1-2 类', '0', '1 类', '5x 广度'],
    ['知识图谱 (动态)', '✓ (千万实体)', '静态', '静态', '✗', '✗', '实时性'],
    ['强化学习', '✓', '✗', '✗', '✗', '✗', '普惠化'],
    ['AI 因子挖掘', '✓ (自动)', '✗', '✗', '半自动', '✗', 'Alpha 加速'],
    ['监管沙盒', '✓ (试点)', '✗', '✗', '✗', '✗', '本土合规'],
    ['中小私募聚焦', '✓ (5000+)', '✗', '✗', '✗', '✗', '差异化蓝海'],
    ['学术合作', '✓ (5 校)', '✗', '✗', '✗', '✓', '信任壁垒'],
    ['综合评价', '全维度领先', '数据全', '散户强', '金融强', '学术强', '—'],
]
add_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.4), data,
          body_size=11, header_size=12, highlight_rows=[1, 2, 6, 7])

add_footer(slide, 6)

# ============== Slide 7: 商业模式 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '六、商业模式', '5 类收入 × 9 类客户 × 订阅主导')

# 5 类收入
add_text(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(0.4),
         '5 类收入结构 (2030)', size=18, bold=True, color=PRIMARY)
revenues = [
    ('SaaS 订阅', '70%', '55,800 万'),
    ('定制开发', '15%', '12,000 万'),
    ('数据 API', '7%', '5,500 万'),
    ('投研培训', '4%', '3,200 万'),
    ('战略合作', '4%', '3,000 万'),
]
for i, (k, p, v) in enumerate(revenues):
    y = 1.85 + i * 0.7
    add_box(slide, Inches(0.4), Inches(y), Inches(6.0), Inches(0.6), fill=LIGHT, line_color=PRIMARY)
    add_text(slide, Inches(0.5), Inches(y + 0.15), Inches(2.5), Inches(0.4),
             k, size=13, bold=True, color=PRIMARY)
    add_text(slide, Inches(3.0), Inches(y + 0.15), Inches(1.0), Inches(0.4),
             p, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.0), Inches(y + 0.15), Inches(2.0), Inches(0.4),
             v, size=12, color=BLACK, align=PP_ALIGN.RIGHT)

# 9 类定价 (右侧)
add_text(slide, Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.4),
         '9 类客户定价 (年付)', size=18, bold=True, color=PRIMARY)
pricing = [
    ('中小私募', '50 万', '银行理财子', '200 万'),
    ('券商资管', '150 万', '信托公司', '100 万'),
    ('保险资管', '300 万', '公募基金', '500 万'),
    ('高校研究所', '20 万', '上市公司 IR', '30 万'),
    ('战略合作', '100 万', '加权 ARPU', '79 万'),
]
for i, (k1, v1, k2, v2) in enumerate(pricing):
    y = 1.85 + i * 0.55
    add_box(slide, Inches(7.0), Inches(y), Inches(2.9), Inches(0.5), fill=WHITE, line_color=SECONDARY)
    add_text(slide, Inches(7.1), Inches(y + 0.1), Inches(1.7), Inches(0.3),
             k1, size=10, color=BLACK)
    add_text(slide, Inches(8.8), Inches(y + 0.1), Inches(1.0), Inches(0.3),
             v1, size=10, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    add_box(slide, Inches(10.0), Inches(y), Inches(2.9), Inches(0.5), fill=WHITE, line_color=SECONDARY)
    add_text(slide, Inches(10.1), Inches(y + 0.1), Inches(1.7), Inches(0.3),
             k2, size=10, color=BLACK)
    add_text(slide, Inches(11.8), Inches(y + 0.1), Inches(1.0), Inches(0.3),
             v2, size=10, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# 销售策略
add_text(slide, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.4),
         '4 大销售策略', size=16, bold=True, color=PRIMARY)
sales = [
    ('直销 80 人', '银行/保险/公募/大型私募'),
    ('渠道联盟', '券商/银行/律所/咨询'),
    ('学术渠道', '清华/北大/上财实验室'),
    ('内容营销', '行业报告/白皮书/技术博客'),
]
for i, (k, v) in enumerate(sales):
    x = 0.4 + i * 3.2
    add_box(slide, Inches(x), Inches(6.2), Inches(3.0), Inches(0.7), fill=LIGHT)
    add_text(slide, Inches(x + 0.1), Inches(6.3), Inches(2.8), Inches(0.3),
             k, size=11, bold=True, color=PRIMARY)
    add_text(slide, Inches(x + 0.1), Inches(6.55), Inches(2.8), Inches(0.3),
             v, size=9, color=BLACK)

add_footer(slide, 7)

# ============== Slide 8: 团队介绍 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '七、团队介绍', '4 创始真名 + 280 团队 (2030) + 5 所高校合作')

# 4 创始
add_text(slide, Inches(0.4), Inches(1.3), Inches(8.0), Inches(0.4),
         '4 位创始团队 (真实身份)', size=18, bold=True, color=PRIMARY)

founders = [
    ('冯亦根', 'CEO', '慧点资本创始人\n10+ 年金融科技投资\nCFA', PRIMARY),
    ('薛永再', 'COO/合规', '杭州永字资管法人\n15+ 年资管行业\n推荐单位代表', ACCENT),
    ('黄成选', 'CTO', '慧点资本量化负责人\nAI 算法专家\n5+ 年量化研究', SECONDARY),
    ('冯思涵', '数据/产品', '慧点资本数据科学家\n金融工程硕士\n3+ 年另类数据', HIGHLIGHT),
]
for i, (name, role, bg, color) in enumerate(founders):
    x = 0.4 + i * 3.2
    add_box(slide, Inches(x), Inches(1.8), Inches(3.0), Inches(2.5),
            fill=WHITE, line_color=color)
    add_box(slide, Inches(x), Inches(1.8), Inches(3.0), Inches(0.6), fill=color)
    add_text(slide, Inches(x), Inches(1.85), Inches(3.0), Inches(0.3),
             name, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(2.15), Inches(3.0), Inches(0.3),
             role, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.15), Inches(2.5), Inches(2.7), Inches(1.7),
             bg, size=10, color=BLACK, align=PP_ALIGN.CENTER)

# 团队扩张
add_text(slide, Inches(0.4), Inches(4.5), Inches(8.0), Inches(0.4),
         '团队扩张规划 (2026-2030)', size=16, bold=True, color=PRIMARY)
team_data = [
    ['部门', '2026', '2027', '2028', '2029', '2030'],
    ['研发 (AI+平台)', '12', '25', '50', '80', '100'],
    ['销售+BD', '3', '10', '25', '50', '80'],
    ['客户成功', '1', '4', '10', '20', '35'],
    ['数据 (采编+ETL)', '3', '8', '15', '25', '35'],
    ['市场+品牌', '1', '3', '5', '10', '15'],
    ['管理', '2', '4', '8', '12', '15'],
    ['合计', '22', '54', '113', '197', '280'],
]
add_table(slide, Inches(0.4), Inches(5.0), Inches(8.5), Inches(2.0), team_data,
          body_size=10, header_size=11, highlight_rows=[7])

# 右侧: 学术合作
add_text(slide, Inches(9.2), Inches(4.5), Inches(3.8), Inches(0.4),
         '学术合作 (5 所)', size=16, bold=True, color=PRIMARY)
academic = [
    '清华五道口金融学院',
    '上财金融科技研究院',
    '北大数字金融研究中心',
    '复旦经济学院',
    '交大/浙大 (校招)',
]
for i, a in enumerate(academic):
    add_box(slide, Inches(9.2), Inches(5.0 + i*0.42), Inches(3.8), Inches(0.38), fill=LIGHT)
    add_text(slide, Inches(9.3), Inches(5.05 + i*0.42), Inches(3.6), Inches(0.3),
             f'✓ {a}', size=11, color=BLACK)

add_footer(slide, 8)

# ============== Slide 9: 5 年里程碑 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '八、5 年里程碑', '12/24/36/48/60 月 6 节点 + 5 轮融资')

milestones = [
    ('M6', 'MVP+5 客户', '26 人', '5 亿', '天使 500 万', '立项+团队组建'),
    ('M12', '产品 2.0+30 客户', '40 人', '3 亿', 'Pre-A 3000 万', 'A 轮准备'),
    ('M24', '100 客户+销售成型', '80 人', '8 亿', 'A 轮 1 亿', '客户成功团队'),
    ('M36', '200 客户+Top 3', '130 人', '20 亿', 'B 轮 3 亿', '海外 5 家'),
    ('M48', '400 客户+壁垒', '200 人', '35 亿', 'Pre-IPO', '专精特新'),
    ('M60', '620 客户+龙头', '280 人', '40-50 亿', '拟 IPO', '启动 IPO 申报'),
]
y0 = 1.4
for i, (m, event, team, val, fund, key) in enumerate(milestones):
    x = 0.4 + i * 2.15
    color = ACCENT if i in [2, 5] else PRIMARY
    add_box(slide, Inches(x), Inches(y0), Inches(2.0), Inches(0.6), fill=color)
    add_text(slide, Inches(x), Inches(y0 + 0.15), Inches(2.0), Inches(0.3),
             m, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(x), Inches(y0 + 0.6), Inches(2.0), Inches(0.7), fill=LIGHT)
    add_text(slide, Inches(x), Inches(y0 + 0.75), Inches(2.0), Inches(0.5),
             event, size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(x), Inches(y0 + 1.3), Inches(2.0), Inches(1.5), fill=WHITE, line_color=color)
    add_text(slide, Inches(x + 0.1), Inches(y0 + 1.4), Inches(1.8), Inches(0.3),
             f'团队: {team}', size=9, color=BLACK)
    add_text(slide, Inches(x + 0.1), Inches(y0 + 1.7), Inches(1.8), Inches(0.3),
             f'估值: {val}', size=9, color=ACCENT, bold=True)
    add_text(slide, Inches(x + 0.1), Inches(y0 + 2.0), Inches(1.8), Inches(0.3),
             f'融资: {fund}', size=9, color=SECONDARY, bold=True)
    add_text(slide, Inches(x + 0.1), Inches(y0 + 2.3), Inches(1.8), Inches(0.5),
             f'关键: {key}', size=8, color=GRAY)

# 底部: KPI
add_text(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.4),
         '董事会月报 10 大 KPI', size=14, bold=True, color=PRIMARY)
kpis = 'MRR/ARR · NRR>110% · CAC 回收<12月 · CSAT>4.5 · NPS>50 · 模型 Sharpe>1.0 · eNPS>40 · 2027Q4 现金跑正 · 年度合规审计 · 数据安全 0 起'
add_box(slide, Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.6), fill=LIGHT)
add_text(slide, Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.4),
         kpis, size=10, color=BLACK, align=PP_ALIGN.CENTER)

add_footer(slide, 9)

# ============== Slide 10: 财务预测 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '九、财务预测', '5 年营收 CAGR 147%, 净利率稳态 25%')

fin_data = [
    ['项目', '2026', '2027', '2028', '2029', '2030', '5年合计'],
    ['客户数 (家)', '30', '85', '200', '380', '620', '1,315'],
    ['ARPU (万/家)', '50', '65', '75', '82', '90', '79 加权'],
    ['营业总收入 (万)', '2,000', '7,300', '20,100', '42,960', '79,500', '151,860'],
    ['同比增长', '-', '265%', '175%', '114%', '85%', 'CAGR 147%'],
    ['毛利润 (万)', '760', '4,250', '14,300', '32,410', '62,000', '113,720'],
    ['毛利率', '38%', '58%', '71%', '75%', '78%', '75% 稳态'],
    ['净利润 (万)', '-1,840', '-800', '2,925', '9,310', '19,800', '29,395'],
    ['净利率', '-92%', '-11%', '15%', '22%', '25%', '22% 加权'],
    ['团队 (人)', '26', '58', '117', '201', '284', '-'],
    ['LTV/CAC', '6.7', '11.8', '18.0', '26.7', '40.0', '>>> 3 金标准'],
]
add_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.5), fin_data,
          body_size=12, header_size=12, highlight_rows=[3, 7])

# 现金流跑正
add_text(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.4),
         '现金流跑正路径', size=14, bold=True, color=PRIMARY)
cashflow = [
    ('2026 启动', '天使+Pre-A 资金支撑', PRIMARY),
    ('2027 Q4', 'OCF 跑正', SECONDARY),
    ('2028 转盈', '净利 +2925 万', ACCENT),
    ('2029 加速', '净利 +9310 万', HIGHLIGHT),
    ('2030 成熟', '净利 +1.98 亿', WARNING),
]
for i, (k, v, c) in enumerate(cashflow):
    x = 0.4 + i * 2.55
    add_box(slide, Inches(x), Inches(6.6), Inches(2.4), Inches(0.5), fill=c)
    add_text(slide, Inches(x), Inches(6.65), Inches(2.4), Inches(0.4),
             f'{k}: {v}', size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_footer(slide, 10)

# ============== Slide 11: 5 场景敏感性 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '十、5 场景敏感性', '熊/保守/基/乐/牛 5 场景, IRR 5%-45%')

scen_data = [
    ['场景', '客户数', '营收 (亿)', '净利率', '净利润 (亿)', '估值 (亿)', 'IRR'],
    ['极端熊市', '310', '3.5', '10%', '0.35', '12', '5%'],
    ['保守', '434', '5.5', '15%', '0.83', '22', '15%'],
    ['基准', '620', '7.95', '25%', '1.98', '40', '28%'],
    ['乐观', '806', '10.3', '30%', '3.09', '60', '38%'],
    ['极端牛市', '992', '12.7', '33%', '4.19', '85', '45%'],
    ['概率加权期望', '-', '-', '-', '-', '39.9', '27%'],
]
add_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.0), scen_data,
          body_size=12, header_size=12, highlight_rows=[3, 6])

# 关键假设
add_text(slide, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.4),
         '关键假设与概率分布', size=16, bold=True, color=PRIMARY)
assum = [
    ('基准', '客户 CAGR 113%, ARPU 12%'),
    ('极端熊市', '客户获取 -50%, 单价 -20%'),
    ('极端牛市', '大客户中标 +30%, 单价 +20%'),
    ('概率分布', '基准 50% + 乐观 25% + 保守 20% + 极端 5%'),
]
for i, (k, v) in enumerate(assum):
    x = 0.4 + (i % 2) * 6.3
    y = 5.1 + (i // 2) * 0.7
    add_box(slide, Inches(x), Inches(y), Inches(6.1), Inches(0.6), fill=LIGHT, line_color=PRIMARY)
    add_text(slide, Inches(x + 0.1), Inches(y + 0.15), Inches(1.5), Inches(0.3),
             k, size=12, bold=True, color=PRIMARY)
    add_text(slide, Inches(x + 1.6), Inches(y + 0.15), Inches(4.4), Inches(0.3),
             v, size=11, color=BLACK)

# 概率加权
add_box(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5), fill=ACCENT)
add_text(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4),
         '期望 IRR = 27% | 期望估值 = 39.9 亿 | 极重压力下 IRR 仍 3% (不亏损)',
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_footer(slide, 11)

# ============== Slide 12: 融资与估值 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '十一、融资与估值', '5 轮 4.35 亿, 估值 5 亿→40 亿, 3 估值方法三角验证')

# 融资轮次
rounds = [
    ('天使', '6 月', '500 万', '5 亿', '10%', '25x', '慧点资本+推荐单位'),
    ('Pre-A', '12 月', '3,000 万', '3 亿', '10%', '20x', '一线 VC'),
    ('A', '24 月', '1 亿', '8 亿', '12.5%', '11x', '一线 VC+产业资本'),
    ('B', '36 月', '3 亿', '20 亿', '15%', '10x', '顶级 PE+战投'),
    ('Pre-IPO', '54 月', '0', '50 亿', '0%', '8x', '二级市场'),
]
y0 = 1.4
for i, (r, t, amt, val, dil, ps, ldr) in enumerate(rounds):
    color = ACCENT if i in [2, 4] else PRIMARY
    add_box(slide, Inches(0.4 + i*2.55), Inches(y0), Inches(2.4), Inches(0.5), fill=color)
    add_text(slide, Inches(0.4 + i*2.55), Inches(y0 + 0.1), Inches(2.4), Inches(0.3),
             r, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(0.4 + i*2.55), Inches(y0 + 0.5), Inches(2.4), Inches(1.5), fill=WHITE, line_color=color)
    add_text(slide, Inches(0.4 + i*2.55 + 0.1), Inches(y0 + 0.6), Inches(2.2), Inches(0.3),
             f'时点: {t}', size=10, color=BLACK)
    add_text(slide, Inches(0.4 + i*2.55 + 0.1), Inches(y0 + 0.9), Inches(2.2), Inches(0.3),
             f'金额: {amt}', size=10, bold=True, color=PRIMARY)
    add_text(slide, Inches(0.4 + i*2.55 + 0.1), Inches(y0 + 1.2), Inches(2.2), Inches(0.3),
             f'估值: {val}', size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(0.4 + i*2.55 + 0.1), Inches(y0 + 1.5), Inches(2.2), Inches(0.3),
             f'稀释: {dil} · {ps}', size=9, color=GRAY)
    add_text(slide, Inches(0.4 + i*2.55 + 0.1), Inches(y0 + 1.8), Inches(2.2), Inches(0.2),
             f'{ldr}', size=8, color=GRAY)

# 3 估值方法
add_text(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.4),
         '估值三角验证 (40 亿, 2030)', size=16, bold=True, color=PRIMARY)
val_data = [
    ['方法', '范围', '依据'],
    ['DCF 估值', '35-40 亿', 'WACC 12%, 永续增长 3%'],
    ['可比 P/S', '12-15x', 'Wind/同花顺/恒生 行业 2026 中位数'],
    ['可比 P/E', '45-60 亿', '成熟期 30x (假设 2030 净利 2 亿)'],
    ['风险溢价', '+5-8%', '资管科技 β 1.3 (相对沪深 300)'],
    ['综合估值', '30-50 亿', '40 亿为基准 (2030 末)'],
]
add_table(slide, Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.4), val_data,
          body_size=12, header_size=12, highlight_rows=[5])

add_footer(slide, 12)

# ============== Slide 13: 风险与应对 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '十二、风险与应对', '8 大风险 + 4 级压力测试, 极重压力下 IRR 仍 3%')

risk_data = [
    ['风险类型', '概率', '营收影响', 'IRR 影响', '应对措施', '责任人'],
    ['AI 技术迭代滞后', '20%', '-15%', '-8pp', '持续 R&D 15% + DeepSeek/Qwen 合作', '黄成选'],
    ['大客户流失', '15%', '-10%', '-5pp', '客户成功 2x + 3 年长约', '冯思涵'],
    ['监管收紧', '30%', '-8%', '-3pp', '法务+合规双保险 + 行业协会', '黄成选'],
    ['融资环境恶化', '25%', '0', '-5pp', '18 月现金跑道 + 多元融资', '薛永再'],
    ['核心团队流失', '15%', '-5%', '-4pp', '期权池 15% + 4 年 vesting', '冯亦根'],
    ['竞品价格战', '40%', '-20%', '-6pp', '差异化 + 成本领先 30%', '冯亦根'],
    ['数据源断供', '10%', '-25%', '-12pp', '多源备份 + 自建爬虫', '黄成选'],
    ['宏观经济衰退', '25%', '-15%', '-7pp', '9 类客户分散 + 抗周期', '冯亦根'],
]
add_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.5), risk_data,
          body_size=10, header_size=11)

# 压力测试
add_text(slide, Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.4),
         '4 级压力测试 (2030 关键指标)', size=16, bold=True, color=PRIMARY)
stress = [
    ('轻度', '客户 590 · 净利 1.6 亿 · 估值 32 亿 · IRR 23%', WARNING),
    ('中度', '客户 530 · 净利 1.2 亿 · 估值 24 亿 · IRR 17%', SECONDARY),
    ('重度', '客户 450 · 净利 0.8 亿 · 估值 16 亿 · IRR 10%', ACCENT),
    ('极重', '客户 350 · 净利 0.3 亿 · 估值 9 亿 · IRR 3%', PRIMARY),
]
for i, (k, v, c) in enumerate(stress):
    y = 5.6 + i * 0.32
    add_box(slide, Inches(0.4), Inches(y), Inches(1.0), Inches(0.3), fill=c)
    add_text(slide, Inches(0.4), Inches(y + 0.05), Inches(1.0), Inches(0.2),
             k, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), Inches(y + 0.05), Inches(11.4), Inches(0.2),
             v, size=10, color=BLACK)

add_footer(slide, 13)

# ============== Slide 14: 退出路径 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '十三、退出路径', '60 月拟 IPO, 4 路径设计 (科创/北交/战投/港股)')

paths = [
    ('A: 科创板 IPO', '60 月启动申报 · 66 月上市\n门槛: 市值≥30 亿, 营收≥5 亿, 净利≥1 亿\n主题: 金融科技 (上海/北京/深圳 试点)\n适配性: ★★★★★', PRIMARY),
    ('B: 北交所 IPO', '54 月启动申报 · 60 月上市\n门槛: 市值≥10 亿, 营收≥2 亿\n主题: 专精特新, 中小机构友好\n适配性: ★★★★', SECONDARY),
    ('C: 战投并购', '48 月可启动 (兜底)\n买方: 头部券商/互联网巨头/资管科技\n估值: 8-10x P/S = 30-40 亿\n适配性: ★★★', ACCENT),
    ('D: 港股 IPO', '60 月可启动 (海外)\n服务: 海外中资机构 + 一带一路\n适配性: ★★★', HIGHLIGHT),
]
y0 = 1.4
for i, (name, desc, color) in enumerate(paths):
    y = y0 + i * 1.4
    add_box(slide, Inches(0.4), Inches(y), Inches(2.0), Inches(1.2), fill=color)
    add_text(slide, Inches(0.4), Inches(y + 0.4), Inches(2.0), Inches(0.4),
             name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(2.4), Inches(y), Inches(10.5), Inches(1.2), fill=WHITE, line_color=color)
    add_text(slide, Inches(2.5), Inches(y + 0.1), Inches(10.4), Inches(1.0),
             desc, size=11, color=BLACK)

# 退出时间表
add_text(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
         '主路径: A (科创板 IPO) 适配性最高, 60 月启动 · 66 月上市 · 估值 40-50 亿',
         size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide, 14)

# ============== Slide 15: 核心竞争优势 ==============
slide = prs.slides.add_slide(BLANK)
add_title_bar(slide, '十四、5 大核心竞争优势', '差异化蓝海 + AI 投研智能体 + 学术背书 + 合规先行 + 团队真实')

adv = [
    ('差异化蓝海', '中小私募 5000+ 家\n被 Wind/同花顺 忽视\n客单价 30-80 万\n✓ 5 年累计 500 家', PRIMARY),
    ('AI 投研智能体', '业内 90% 不具备\n12-18 月技术窗口\nDeepSeek/Qwen 微调\n✓ 自动因子挖掘', ACCENT),
    ('学术合作背书', '清华/北大/上财/复旦/交大\n监管沙盒试点\n学术顾问 5+ 位\n✓ 信任壁垒', SECONDARY),
    ('合规先行', '数据/金融/AI 三重合规\n法务+合规双保险\n推荐单位借调 6 月\n✓ 本土合规', HIGHLIGHT),
    ('团队真实可查', '4 创始慧点资本实际团队\n推荐单位法定代表人\n资管行业 15 年资源\n✓ 无失实描述', WARNING),
]
y0 = 1.4
for i, (name, desc, color) in enumerate(adv):
    x = 0.4 + (i % 5) * 2.55
    y = y0
    add_box(slide, Inches(x), Inches(y), Inches(2.4), Inches(0.5), fill=color)
    add_text(slide, Inches(x), Inches(y + 0.1), Inches(2.4), Inches(0.3),
             f'#{i+1} {name}', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(x), Inches(y + 0.5), Inches(2.4), Inches(2.0), fill=WHITE, line_color=color)
    add_text(slide, Inches(x + 0.1), Inches(y + 0.6), Inches(2.2), Inches(1.9),
             desc, size=10, color=BLACK)

# 4 大数据
add_text(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.4),
         '评委核心关注响应', size=18, bold=True, color=PRIMARY)
data_pts = [
    ('盈利能力', '净利率 25% 行业 Top 20%', '强于恒生/通联'),
    ('增长持续性', '5 年 CAGR 147%, 9 类客户分群', '多行业抗周期'),
    ('现金流', '2027 Q4 OCF 跑正, 不依赖持续融资', '财务健康'),
    ('估值合理性', '40 亿 5x P/S, DCF/可比/风险溢价三角验证', '多方法交叉'),
    ('退出路径', '60 月拟 IPO (科创板/北交所)', '明确时间表'),
]
for i, (k, v, note) in enumerate(data_pts):
    y = 4.5 + i * 0.45
    add_box(slide, Inches(0.4), Inches(y), Inches(2.5), Inches(0.4), fill=PRIMARY)
    add_text(slide, Inches(0.4), Inches(y + 0.05), Inches(2.5), Inches(0.3),
             k, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.0), Inches(y + 0.05), Inches(7.0), Inches(0.3),
             v, size=11, color=BLACK)
    add_text(slide, Inches(10.0), Inches(y + 0.05), Inches(3.0), Inches(0.3),
             f'✓ {note}', size=10, color=ACCENT, bold=True)

add_footer(slide, 15)

# ============== Slide 16: 总结与愿景 ==============
slide = prs.slides.add_slide(BLANK)
# 背景
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.line.fill.background()
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY

# 副标题条
sub = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), SLIDE_W, Inches(0.1))
sub.line.fill.background()
sub.fill.solid()
sub.fill.fore_color.rgb = ACCENT

# 主标题
add_text(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.5),
         '十五、总结与愿景', size=24, bold=True, color=WHITE)

# 一页纸摘要
add_text(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
         '让量化投资更智能, 让另类数据更普惠, 让资管科技更本土化', size=18, color=WHITE, align=PP_ALIGN.CENTER)

# 关键数字
key_nums = [
    ('620 家', '客户 (2030)'),
    ('7.95 亿', '营收 (2030)'),
    ('1.98 亿', '净利润 (2030)'),
    ('25%', '净利率 (稳态)'),
    ('40 亿', '估值 (2030)'),
    ('28%', 'IRR (5 年)'),
]
y0 = 3.2
for i, (n, l) in enumerate(key_nums):
    x = 0.5 + (i % 6) * 2.05
    add_box(slide, Inches(x), Inches(y0), Inches(1.9), Inches(1.0), fill=WHITE)
    add_text(slide, Inches(x), Inches(y0 + 0.15), Inches(1.9), Inches(0.5),
             n, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(y0 + 0.6), Inches(1.9), Inches(0.4),
             l, size=10, color=PRIMARY, align=PP_ALIGN.CENTER)

# 5 大交付
add_text(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         '完整提交包 (16+ 交付物)', size=16, bold=True, color=WHITE)
delivs = [
    '✓ BP V2.0 (14 页 PDF)',
    '✓ 财务模型 V3 (13 工作表)',
    '✓ 技术白皮书 (19 页 PDF)',
    '✓ 3 分钟 Demo 视频',
    '✓ 3 条技术短视频',
    '✓ 8 分钟路演视频',
    '✓ Demo 原型 (Streamlit)',
    '✓ 客户证据包 (6 文件)',
    '✓ T13 财务图表 (V3)',
    '✓ 行业研究 (T01)',
    '✓ T11 验收报告',
    '✓ T12-T16 验收报告',
]
for i, d in enumerate(delivs):
    x = 0.5 + (i % 4) * 3.15
    y = 5.0 + (i // 4) * 0.4
    add_text(slide, Inches(x), Inches(y), Inches(3.0), Inches(0.3),
             d, size=11, color=WHITE)

# 联系方式
add_text(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
         '慧点资本 (InsightQuant) 量化研究部  |  杭州永字资产管理有限公司 (推荐单位)',
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         '项目编号: 2026FINTECH-FINT-0093  |  FinTech@外滩金融科技大赛 · 资管科技赛道  |  2026 年 6 月',
         size=12, color=WHITE, align=PP_ALIGN.CENTER)

# 保存
output_path = 'D:/shFintech/QuantInsight_Pro_Pitch_Deck_V2.pptx'
prs.save(output_path)
size = os.path.getsize(output_path) / 1024
print(f'[OK] PPT V2: {output_path}')
print(f'[OK] 大小: {size:.0f} KB')
print(f'[OK] 幻灯片: {len(prs.slides)}')
